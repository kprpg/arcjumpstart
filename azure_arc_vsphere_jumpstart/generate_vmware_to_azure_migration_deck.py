from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).with_name("VMware_to_Azure_Migration_Deck.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

COLORS = {
    "navy": RGBColor(11, 33, 68),
    "blue": RGBColor(0, 120, 212),
    "teal": RGBColor(0, 153, 153),
    "orange": RGBColor(255, 140, 0),
    "cream": RGBColor(247, 249, 252),
    "slate": RGBColor(58, 74, 91),
    "text": RGBColor(28, 37, 44),
    "muted": RGBColor(94, 107, 122),
    "white": RGBColor(255, 255, 255),
    "line": RGBColor(214, 221, 230),
}

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"


def add_background(slide, accent_color=COLORS["blue"]):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["cream"]

    header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.42)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS["navy"]
    header.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        Inches(0.42),
        Inches(0.18),
        prs.slide_height - Inches(0.42),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()


def add_footer(slide, source_text):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.35),
        Inches(7.0),
        Inches(12.55),
        Inches(0.015),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()

    box = slide.shapes.add_textbox(
        Inches(0.45), Inches(7.03), Inches(9.8), Inches(0.26)
    )
    frame = box.text_frame
    p = frame.paragraphs[0]
    p.text = source_text
    p.font.name = BODY_FONT
    p.font.size = Pt(8)
    p.font.color.rgb = COLORS["muted"]

    brand = slide.shapes.add_textbox(
        Inches(10.55), Inches(6.99), Inches(2.2), Inches(0.3)
    )
    frame = brand.text_frame
    p = frame.paragraphs[0]
    p.text = "arcjumpstart | Azure Arc vSphere"
    p.font.name = BODY_FONT
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = COLORS["slate"]
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, title, subtitle=None, color=COLORS["navy"]):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.75), Inches(11.6), Inches(0.9))
    frame = box.text_frame
    p = frame.paragraphs[0]
    p.text = title
    p.font.name = TITLE_FONT
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = color

    if subtitle:
        sub = slide.shapes.add_textbox(
            Inches(0.72), Inches(1.55), Inches(11.2), Inches(0.5)
        )
        frame = sub.text_frame
        p = frame.paragraphs[0]
        p.text = subtitle
        p.font.name = BODY_FONT
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS["muted"]


def add_bullets(
    slide, left, top, width, height, bullets, font_size=18, color=COLORS["text"]
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Pt(6)
    frame.margin_right = Pt(6)
    frame.margin_top = Pt(4)
    frame.vertical_anchor = MSO_ANCHOR.TOP

    first = True
    for bullet in bullets:
        p = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        p.text = bullet
        p.level = 0
        p.font.name = BODY_FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.bullet = True


def add_callout(slide, left, top, width, height, title, body, fill):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill

    title_box = slide.shapes.add_textbox(
        left + Inches(0.18), top + Inches(0.12), width - Inches(0.3), Inches(0.35)
    )
    frame = title_box.text_frame
    p = frame.paragraphs[0]
    p.text = title
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]

    body_box = slide.shapes.add_textbox(
        left + Inches(0.18),
        top + Inches(0.5),
        width - Inches(0.3),
        height - Inches(0.6),
    )
    frame = body_box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = body
    p.font.name = BODY_FONT
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS["white"]


def add_metric(slide, left, top, width, height, value, caption, accent):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["white"]
    card.line.color.rgb = COLORS["line"]

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.12), height
    )
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()

    value_box = slide.shapes.add_textbox(
        left + Inches(0.22), top + Inches(0.15), width - Inches(0.3), Inches(0.45)
    )
    frame = value_box.text_frame
    p = frame.paragraphs[0]
    p.text = value
    p.font.name = TITLE_FONT
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS["navy"]

    caption_box = slide.shapes.add_textbox(
        left + Inches(0.22),
        top + Inches(0.7),
        width - Inches(0.3),
        height - Inches(0.8),
    )
    frame = caption_box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = caption
    p.font.name = BODY_FONT
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS["slate"]


def add_section_divider(title, subtitle, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["navy"]

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.8),
        Inches(0.95),
        Inches(0.18),
        Inches(4.8),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(1.3), Inches(1.0), Inches(10.7), Inches(1.2)
    )
    frame = title_box.text_frame
    p = frame.paragraphs[0]
    p.text = title
    p.font.name = TITLE_FONT
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]

    sub_box = slide.shapes.add_textbox(
        Inches(1.32), Inches(2.15), Inches(9.5), Inches(1.0)
    )
    frame = sub_box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = subtitle
    p.font.name = BODY_FONT
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(219, 231, 245)

    add_footer(slide, "Section transition for mixed executive and technical audiences")


def add_two_column_slide(
    title,
    subtitle,
    left_title,
    left_bullets,
    right_title,
    right_bullets,
    footer,
    accent=COLORS["blue"],
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, accent)
    add_title(slide, title, subtitle)

    for x, heading in ((0.75, left_title), (6.65, right_title)):
        label = slide.shapes.add_textbox(
            Inches(x), Inches(2.05), Inches(5.4), Inches(0.35)
        )
        p = label.text_frame.paragraphs[0]
        p.text = heading
        p.font.name = BODY_FONT
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = accent

    left_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.72),
        Inches(2.35),
        Inches(5.55),
        Inches(4.2),
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = COLORS["white"]
    left_panel.line.color.rgb = COLORS["line"]

    right_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.62),
        Inches(2.35),
        Inches(5.55),
        Inches(4.2),
    )
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = COLORS["white"]
    right_panel.line.color.rgb = COLORS["line"]

    add_bullets(
        slide,
        Inches(0.92),
        Inches(2.58),
        Inches(5.1),
        Inches(3.8),
        left_bullets,
        font_size=16,
    )
    add_bullets(
        slide,
        Inches(6.82),
        Inches(2.58),
        Inches(5.1),
        Inches(3.8),
        right_bullets,
        font_size=16,
    )
    add_footer(slide, footer)


def add_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["navy"]

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.7),
        Inches(0.8),
        Inches(0.22),
        Inches(5.8),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["blue"]
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(1.15), Inches(1.0), Inches(9.7), Inches(1.7)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = "VMware to Azure Migration"
    p.font.name = TITLE_FONT
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]

    p = title_box.text_frame.add_paragraph()
    p.text = "Executive and technical decision deck built from the VMware-to-Azure eBook and Microsoft Learn guidance"
    p.font.name = BODY_FONT
    p.font.size = Pt(17)
    p.font.color.rgb = RGBColor(219, 231, 245)
    p.space_before = Pt(8)

    ribbon = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(1.18),
        Inches(3.0),
        Inches(4.0),
        Inches(0.56),
    )
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = COLORS["teal"]
    ribbon.line.fill.background()
    p = slide.shapes.add_textbox(
        Inches(1.42), Inches(3.13), Inches(3.6), Inches(0.25)
    ).text_frame.paragraphs[0]
    p.text = "Business leaders + technical decision makers"
    p.font.name = BODY_FONT
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]

    body_box = slide.shapes.add_textbox(
        Inches(1.18), Inches(3.95), Inches(5.6), Inches(1.6)
    )
    frame = body_box.text_frame
    frame.word_wrap = True
    for idx, bullet in enumerate(
        [
            "Why customers move from VMware to Azure now",
            "When to choose Azure VMware Solution, Azure VMs/PaaS, Azure Arc, or Azure Virtual Desktop",
            "How to structure discovery, landing zones, migration waves, and cost optimization",
        ]
    ):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = bullet
        p.font.name = BODY_FONT
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS["white"]
        p.bullet = True
        p.space_after = Pt(8)

    add_callout(
        slide,
        Inches(8.2),
        Inches(1.3),
        Inches(4.25),
        Inches(1.15),
        "Primary source",
        "The Ultimate Guide to VMware Migration and Modernization on Azure (2025 eBook)",
        COLORS["blue"],
    )
    add_callout(
        slide,
        Inches(8.2),
        Inches(2.7),
        Inches(4.25),
        Inches(2.15),
        "Official Microsoft guidance used",
        "Azure Migrate start-here for VMware, business case guidance, Azure VMware Solution landing zone guidance, and Well-Architected Azure VMware Solution migration content.",
        COLORS["orange"],
    )
    add_footer(
        slide,
        "Sources: VMware-to-Azure eBook (2025) and Microsoft Learn Azure Migrate / Azure VMware Solution guidance",
    )


def add_audience_map_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["teal"])
    add_title(
        slide,
        "How to use this deck",
        "The first section is tuned for business outcomes; the second turns those choices into a migration plan.",
    )

    add_callout(
        slide,
        Inches(0.85),
        Inches(2.0),
        Inches(3.8),
        Inches(2.35),
        "Business leaders",
        "Focus on urgency, target-state options, financial outcomes, risk reduction, and phased investment decisions.",
        COLORS["blue"],
    )
    add_callout(
        slide,
        Inches(4.75),
        Inches(2.0),
        Inches(3.8),
        Inches(2.35),
        "Technical decision makers",
        "Focus on discovery, assessment, target architecture, landing zone preparation, migration tooling, and hybrid operations.",
        COLORS["teal"],
    )
    add_callout(
        slide,
        Inches(8.65),
        Inches(2.0),
        Inches(3.8),
        Inches(2.35),
        "Shared decisions",
        "Decide workload placement, sequencing, governance, networking, identity, and cost controls before execution begins.",
        COLORS["orange"],
    )

    add_bullets(
        slide,
        Inches(0.95),
        Inches(4.85),
        Inches(11.2),
        Inches(1.5),
        [
            "Use the executive slides to align on strategy and funding.",
            "Use the technical slides to validate feasibility, dependencies, and migration waves.",
            "Use the final action plan as the handoff into Azure Migrate and landing zone work.",
        ],
        font_size=17,
    )
    add_footer(
        slide,
        "Deck structure derived from the eBook chapters plus Microsoft Learn migration workflow guidance",
    )


def add_exec_summary_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["blue"])
    add_title(
        slide,
        "Executive summary",
        "Azure supports more than one VMware migration lane; most enterprises use a mix of them.",
    )

    add_metric(
        slide,
        Inches(0.8),
        Inches(2.0),
        Inches(2.75),
        Inches(1.45),
        "3 paths",
        "Azure VMware Solution, Azure VMs/PaaS, and hybrid extension with Azure Arc.",
        COLORS["blue"],
    )
    add_metric(
        slide,
        Inches(3.75),
        Inches(2.0),
        Inches(2.75),
        Inches(1.45),
        "298% ROI",
        "Composite ROI reported in the Forrester study cited by the eBook for Azure VMware Solution.",
        COLORS["teal"],
    )
    add_metric(
        slide,
        Inches(6.7),
        Inches(2.0),
        Inches(2.75),
        Inches(1.45),
        "$5.1M",
        "Reduced datacenter TCO cited in the eBook's featured value proof points.",
        COLORS["orange"],
    )
    add_metric(
        slide,
        Inches(9.65),
        Inches(2.0),
        Inches(2.75),
        Inches(1.45),
        "60+ regions",
        "Global Azure footprint supports resilience, business continuity, and capacity planning.",
        COLORS["blue"],
    )

    add_bullets(
        slide,
        Inches(0.88),
        Inches(4.0),
        Inches(11.2),
        Inches(2.2),
        [
            "Choose Azure VMware Solution when speed, low disruption, and VMware operating symmetry matter most.",
            "Choose Azure IaaS/PaaS when you want to reduce VMware dependency and modernize apps, databases, or VDI over time.",
            "Keep on-premises vSphere connected with Azure Arc when some workloads must stay put but still need Azure governance, security, and lifecycle controls.",
        ],
        font_size=18,
    )
    add_footer(
        slide,
        "Sources: eBook pages 4-6, 20-21; Microsoft Learn business case and Azure VMware Solution overview guidance",
    )


def add_business_case_slide():
    add_two_column_slide(
        "Why act now",
        "The eBook frames migration as a modernization decision, not just a hosting change.",
        "Business pressure",
        [
            "Aging infrastructure, security exposure, and datacenter cost pressure make status quo harder to justify.",
            "AI readiness increases demand for scalable infrastructure, data services, and cloud-adjacent innovation.",
            "Leadership teams want faster regional expansion, merger support, and resilience without new hardware cycles.",
        ],
        "Azure response",
        [
            "Cloud agility with elastic capacity and global reach.",
            "Code-to-cloud security through Defender for Cloud, Azure Firewall, and policy-driven governance.",
            "Licensing and commercial levers such as Azure Hybrid Benefit, reserved capacity, and Extended Security Updates in Azure.",
        ],
        "Sources: eBook pages 3-6 and 20; Microsoft Learn business case overview",
        accent=COLORS["orange"],
    )


def add_choice_matrix_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["teal"])
    add_title(
        slide,
        "Choose the right migration lane",
        "Microsoft guidance and the eBook both emphasize that there is no single default target for every VMware workload.",
    )

    headers = ["Path", "Best fit", "What you gain", "Tradeoff"]
    rows = [
        [
            "Azure VMware Solution",
            "Fast datacenter exit, minimal refactoring, preserve VMware operations",
            "vSphere symmetry, HCX/vMotion paths, Azure service adjacency",
            "Higher VMware affinity remains in target state",
        ],
        [
            "Azure VMs + PaaS",
            "Intentional modernization of apps, databases, and operations",
            "No VMware license dependency, native Azure services, easier long-term optimization",
            "Requires more planning, landing zone maturity, and app change",
        ],
        [
            "Azure Arc-enabled VMware vSphere",
            "Workloads staying on-premises or moving in phases",
            "Azure governance, RBAC, inventory, policy, and lifecycle operations over vSphere",
            "Not a full migration target by itself",
        ],
        [
            "Azure Virtual Desktop / Horizon Cloud on Azure",
            "VDI transformation or capacity expansion",
            "Cloud-hosted desktops, right-sizing, GPU options, simpler scaling",
            "Separate end-user compute workstream",
        ],
    ]

    table_x = Inches(0.7)
    table_y = Inches(2.0)
    col_widths = [Inches(2.3), Inches(2.6), Inches(3.2), Inches(3.5)]
    row_height = Inches(0.88)

    current_x = table_x
    for idx, header in enumerate(headers):
        cell = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            current_x,
            table_y,
            col_widths[idx],
            Inches(0.55),
        )
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["navy"]
        cell.line.color.rgb = COLORS["navy"]
        p = slide.shapes.add_textbox(
            current_x + Inches(0.08),
            table_y + Inches(0.12),
            col_widths[idx] - Inches(0.16),
            Inches(0.28),
        ).text_frame.paragraphs[0]
        p.text = header
        p.font.name = BODY_FONT
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]
        current_x += col_widths[idx]

    for row_idx, row in enumerate(rows):
        current_x = table_x
        y = table_y + Inches(0.62) + row_idx * row_height
        for col_idx, value in enumerate(row):
            cell = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                current_x,
                y,
                col_widths[col_idx],
                row_height,
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                COLORS["white"] if row_idx % 2 == 0 else RGBColor(242, 246, 250)
            )
            cell.line.color.rgb = COLORS["line"]
            box = slide.shapes.add_textbox(
                current_x + Inches(0.08),
                y + Inches(0.08),
                col_widths[col_idx] - Inches(0.16),
                row_height - Inches(0.1),
            )
            frame = box.text_frame
            frame.word_wrap = True
            p = frame.paragraphs[0]
            p.text = value
            p.font.name = BODY_FONT
            p.font.size = Pt(10.5)
            p.font.color.rgb = COLORS["text"]
            current_x += col_widths[col_idx]

    add_footer(
        slide,
        "Sources: eBook pages 4, 10, 16, 18-19; Microsoft Learn start-here and AVS workload guidance",
    )


def add_discovery_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["blue"])
    add_title(
        slide,
        "Technical migration workflow",
        "Microsoft Learn's VMware migration map aligns closely with the eBook: discover, assess, build the case, prepare the landing zone, then migrate.",
    )

    steps = [
        (
            "1. Discover",
            "Use Azure Migrate appliance discovery or RVTools XLSX import to inventory servers.",
        ),
        (
            "2. Map dependencies",
            "Use agentless or agent-based dependency analysis to define application groups and wave boundaries.",
        ),
        (
            "3. Build the business case",
            "Compare Azure VMs, AVS, and modernization options with TCO and readiness data.",
        ),
        (
            "4. Assess readiness",
            "Identify blockers, sizing recommendations, cost estimates, and conditional readiness.",
        ),
        (
            "5. Migrate and validate",
            "Use HCX/vMotion for AVS or Azure Migrate replication for Azure VMs, then validate and optimize.",
        ),
    ]

    left = Inches(0.88)
    top = Inches(2.1)
    width = Inches(2.2)
    gap = Inches(0.18)
    for idx, (heading, body) in enumerate(steps):
        x = left + idx * (width + gap)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, top, width, Inches(3.4)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS["white"]
        card.line.color.rgb = COLORS["line"]

        circle = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            x + Inches(0.12),
            top + Inches(0.12),
            Inches(0.42),
            Inches(0.42),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS["blue"] if idx < 3 else COLORS["teal"]
        circle.line.fill.background()
        p = slide.shapes.add_textbox(
            x + Inches(0.24), top + Inches(0.2), Inches(0.18), Inches(0.18)
        ).text_frame.paragraphs[0]
        p.text = str(idx + 1)
        p.font.name = BODY_FONT
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]
        p.alignment = PP_ALIGN.CENTER

        heading_box = slide.shapes.add_textbox(
            x + Inches(0.18), top + Inches(0.7), width - Inches(0.36), Inches(0.45)
        )
        p = heading_box.text_frame.paragraphs[0]
        p.text = heading
        p.font.name = BODY_FONT
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS["navy"]

        body_box = slide.shapes.add_textbox(
            x + Inches(0.18), top + Inches(1.12), width - Inches(0.36), Inches(2.0)
        )
        frame = body_box.text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.text = body
        p.font.name = BODY_FONT
        p.font.size = Pt(11.5)
        p.font.color.rgb = COLORS["slate"]

    add_footer(
        slide,
        "Sources: Microsoft Learn start-here for VMware migration and Azure Migrate planning guidance; eBook pages 14-15",
    )


def add_architecture_slide():
    add_two_column_slide(
        "Target-state architecture decisions",
        "Technical leaders should decide placement by application need, operational model, and modernization appetite.",
        "Keep VMware operating symmetry",
        [
            "Land selected workloads in Azure VMware Solution private clouds when vSphere consistency and fast migration matter most.",
            "Connect on-premises vSphere to Azure through ExpressRoute or VPN.",
            "Attach Azure-native services over time for monitoring, security, backup, data, and application evolution.",
        ],
        "Move toward native Azure platforms",
        [
            "Use Azure Migrate assessment data to identify candidates for Azure VMs, Azure SQL, App Service, AKS, or Azure Virtual Desktop.",
            "Use application groups and dependency maps to split quick wins from harder refactors.",
            "Treat AVS and Azure-native landing zones as complementary, not mutually exclusive.",
        ],
        "Sources: eBook pages 10-13; Microsoft Learn business case overview and AVS workload guidance",
        accent=COLORS["teal"],
    )


def add_landing_zone_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["orange"])
    add_title(
        slide,
        "Azure VMware Solution landing zone essentials",
        "Official guidance is explicit: the network and platform foundations need to be in place before migration waves begin.",
    )

    add_bullets(
        slide,
        Inches(0.9),
        Inches(2.0),
        Inches(5.4),
        Inches(3.8),
        [
            "Choose subscription, resource group, region, and naming strategy early.",
            "Estimate host count and request AVS host quota in the target subscription.",
            "Reserve a non-overlapping /22 CIDR for AVS private cloud management.",
            "Deploy or identify the Azure virtual network and gateway needed for connectivity.",
            "Define NSX-T segments, route propagation, and firewall rules before cutover.",
        ],
        font_size=16,
    )

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.6),
        Inches(2.0),
        Inches(5.5),
        Inches(3.95),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS["white"]
    panel.line.color.rgb = COLORS["line"]

    label = slide.shapes.add_textbox(
        Inches(6.85), Inches(2.24), Inches(4.8), Inches(0.4)
    )
    p = label.text_frame.paragraphs[0]
    p.text = "Shared platform responsibilities"
    p.font.name = BODY_FONT
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLORS["orange"]

    add_bullets(
        slide,
        Inches(6.82),
        Inches(2.72),
        Inches(4.95),
        Inches(2.9),
        [
            "Platform landing zones provide identity, network, policy, security, and monitoring services.",
            "Application landing zones give workload teams a governed place to consume those services.",
            "Treat AVS integration with landing zones as a joint platform-and-workload design exercise.",
        ],
        font_size=15,
    )

    add_footer(
        slide,
        "Sources: Microsoft Learn CAF ready/migrate for Azure VMware Solution and landing-zone integration guidance",
    )


def add_waves_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["blue"])
    add_title(
        slide,
        "Wave planning and tooling",
        "The first migration wave should prove connectivity, operational runbooks, rollback, and validation criteria before scale-out.",
    )

    phases = [
        (
            "Pilot",
            "Low-risk app group, validate network path, identity, monitoring, backup, and rollback.",
        ),
        (
            "Factory wave 1",
            "Move straightforward VMs or dev/test groups to establish throughput and team roles.",
        ),
        (
            "Factory wave 2+",
            "Migrate business-critical groups using established templates, cutover playbooks, and approval gates.",
        ),
        (
            "Optimize",
            "Right-size, attach Azure services, decommission stranded infrastructure, and revisit modernization candidates.",
        ),
    ]

    for idx, (heading, body) in enumerate(phases):
        x = Inches(0.95) + idx * Inches(3.02)
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            x,
            Inches(2.45),
            Inches(2.68),
            Inches(2.25),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS["white"]
        box.line.color.rgb = COLORS["line"]

        marker = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, Inches(2.45), Inches(2.68), Inches(0.18)
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = [
            COLORS["teal"],
            COLORS["blue"],
            COLORS["orange"],
            COLORS["navy"],
        ][idx]
        marker.line.fill.background()

        title_box = slide.shapes.add_textbox(
            x + Inches(0.18), Inches(2.76), Inches(2.25), Inches(0.35)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = heading
        p.font.name = BODY_FONT
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS["navy"]

        body_box = slide.shapes.add_textbox(
            x + Inches(0.18), Inches(3.18), Inches(2.22), Inches(1.25)
        )
        frame = body_box.text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.text = body
        p.font.name = BODY_FONT
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS["slate"]

    add_bullets(
        slide,
        Inches(0.92),
        Inches(5.2),
        Inches(11.2),
        Inches(1.0),
        [
            "Tooling by path: HCX / vMotion for AVS rehosting, Azure Migrate appliance + replication for Azure VMs, Azure Arc for hybrid visibility and control.",
        ],
        font_size=15,
    )
    add_footer(
        slide,
        "Sources: eBook pages 14-16; Microsoft Learn server migration overview and VMware migration start-here guidance",
    )


def add_arc_slide():
    add_two_column_slide(
        "Hybrid operations with Azure Arc-enabled VMware vSphere",
        "Arc is the bridge for phased migration programs and for workloads that must remain in vSphere for longer.",
        "What Arc adds",
        [
            "Inventory and register vSphere resources in Azure at scale.",
            "Apply consistent RBAC, governance, and Azure management workflows.",
            "Create, resize, delete, and power-manage VMware VMs through Azure controls.",
        ],
        "Where it fits",
        [
            "Supports hybrid estates where not every workload moves on day one.",
            "Lets platform teams extend governance before, during, and after migration waves.",
            "Creates a cleaner handoff between infrastructure teams and application owners.",
        ],
        "Sources: eBook pages 16-17 and Azure Arc-enabled VMware vSphere documentation referenced by the eBook",
        accent=COLORS["orange"],
    )


def add_vdi_slide():
    add_two_column_slide(
        "VDI and end-user compute options",
        "The eBook treats VDI as a distinct modernization lane, not merely another VM migration wave.",
        "Azure Virtual Desktop",
        [
            "Use cloud-hosted desktops with centralized Azure management, security, and elastic scaling.",
            "Right-size hosts, shut them down when idle, and add GPU-backed pools for design or engineering workloads.",
            "Good fit when you want Microsoft-managed desktop services and tighter Azure integration.",
        ],
        "VMware Horizon Cloud on Azure",
        [
            "Preserves Horizon-centric administration with minimal Azure expertise required.",
            "Can pair with Azure Virtual Desktop to combine VMware operational familiarity with Azure infrastructure.",
            "Useful when desktop ops teams want a VMware-first management experience during transition.",
        ],
        "Sources: eBook pages 18-19",
        accent=COLORS["teal"],
    )


def add_cost_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["orange"])
    add_title(
        slide,
        "Economics and optimization levers",
        "Use the business case to compare on-premises costs against Azure VMs, Azure VMware Solution, and modernization scenarios.",
    )

    add_metric(
        slide,
        Inches(0.8),
        Inches(1.95),
        Inches(2.75),
        Inches(1.35),
        "AHB",
        "Use Azure Hybrid Benefit for Windows Server and SQL Server to lower compute cost.",
        COLORS["blue"],
    )
    add_metric(
        slide,
        Inches(3.8),
        Inches(1.95),
        Inches(2.75),
        Inches(1.35),
        "Reserved",
        "Lock in pricing for steady-state workloads with reserved capacity where utilization is predictable.",
        COLORS["teal"],
    )
    add_metric(
        slide,
        Inches(6.8),
        Inches(1.95),
        Inches(2.75),
        Inches(1.35),
        "ESU",
        "Use Extended Security Updates in Azure for older Windows Server and SQL Server estates when required.",
        COLORS["orange"],
    )
    add_metric(
        slide,
        Inches(9.8),
        Inches(1.95),
        Inches(2.75),
        Inches(1.35),
        "Right-size",
        "Use Azure Migrate assessment data to resize compute based on actual CPU, memory, and storage usage.",
        COLORS["navy"],
    )

    add_two_column_costs = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.82),
        Inches(3.7),
        Inches(11.35),
        Inches(2.5),
    )
    add_two_column_costs.fill.solid()
    add_two_column_costs.fill.fore_color.rgb = COLORS["white"]
    add_two_column_costs.line.color.rgb = COLORS["line"]

    left_label = slide.shapes.add_textbox(
        Inches(1.05), Inches(4.0), Inches(4.8), Inches(0.35)
    )
    p = left_label.text_frame.paragraphs[0]
    p.text = "Business case viewpoints"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS["orange"]

    add_bullets(
        slide,
        Inches(1.0),
        Inches(4.35),
        Inches(5.0),
        Inches(1.45),
        [
            "Compare on-premises TCO against Azure VMs, AVS, and modernization-to-PaaS scenarios.",
            "Include savings drivers such as license reuse, security updates, and operating model change.",
        ],
        font_size=14,
    )

    right_label = slide.shapes.add_textbox(
        Inches(6.45), Inches(4.0), Inches(4.8), Inches(0.35)
    )
    p = right_label.text_frame.paragraphs[0]
    p.text = "What to avoid"
    p.font.name = BODY_FONT
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS["orange"]

    add_bullets(
        slide,
        Inches(6.4),
        Inches(4.35),
        Inches(5.0),
        Inches(1.45),
        [
            "Do not price migration only as infrastructure replication; account for stranded assets, skill transitions, and decommission timing.",
            "Do not default every workload to the same target if supportability or cost outcomes differ.",
        ],
        font_size=14,
    )
    add_footer(
        slide,
        "Sources: eBook page 20; Microsoft Learn business case and VMware migration assessment guidance",
    )


def add_risks_slide():
    add_two_column_slide(
        "Common risks and mitigation patterns",
        "These are the issues that most often derail VMware-to-Azure programs when they are discovered too late.",
        "Risk areas",
        [
            "Unmapped dependencies create broken application cutovers and hidden downtime.",
            "Overlapping IP ranges or late network design delay AVS connectivity and testing.",
            "Unclear workload placement rules cause rework across AVS, Azure VMs, and modernization teams.",
        ],
        "Mitigations",
        [
            "Use dependency analysis and application grouping before wave scheduling.",
            "Lock landing zone, identity, route, firewall, and CIDR decisions before pilot migration.",
            "Create a placement matrix with business owner sign-off and update it after each assessment cycle.",
        ],
        "Sources: Microsoft Learn start-here and AVS landing-zone guidance, reinforced by eBook pages 14-16",
        accent=COLORS["blue"],
    )


def add_action_plan_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["teal"])
    add_title(
        slide,
        "Recommended 90-day action plan",
        "Use this as the working close for stakeholder alignment and next-step ownership.",
    )

    plan = [
        (
            "Days 1-30",
            [
                "Create or confirm the Azure Migrate project and collect RVTools or appliance-based discovery data.",
                "Segment workloads into AVS candidates, Azure-native modernization candidates, VDI candidates, and stay-on-premises groups.",
            ],
        ),
        (
            "Days 31-60",
            [
                "Build the business case, validate landing zone assumptions, and request quotas for target subscriptions.",
                "Run readiness assessments and choose a pilot wave with rollback criteria.",
            ],
        ),
        (
            "Days 61-90",
            [
                "Execute the pilot migration, validate operations, and refine cutover runbooks.",
                "Approve factory-wave sequencing, decommission plans, and modernization backlog ownership.",
            ],
        ),
    ]

    for idx, (heading, bullets) in enumerate(plan):
        x = Inches(0.95) + idx * Inches(4.05)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            x,
            Inches(2.1),
            Inches(3.55),
            Inches(3.9),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS["white"]
        card.line.color.rgb = COLORS["line"]

        header = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, Inches(2.1), Inches(3.55), Inches(0.22)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = [COLORS["blue"], COLORS["teal"], COLORS["orange"]][
            idx
        ]
        header.line.fill.background()

        title_box = slide.shapes.add_textbox(
            x + Inches(0.2), Inches(2.45), Inches(3.1), Inches(0.35)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = heading
        p.font.name = BODY_FONT
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLORS["navy"]

        add_bullets(
            slide,
            x + Inches(0.12),
            Inches(2.85),
            Inches(3.15),
            Inches(2.7),
            bullets,
            font_size=14,
        )

    add_footer(
        slide,
        "Recommended action plan synthesized from the eBook and Microsoft Learn VMware migration workflow guidance",
    )


def add_sources_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["navy"])

    title_box = slide.shapes.add_textbox(
        Inches(0.72), Inches(0.8), Inches(11.2), Inches(0.8)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = "Source references used in this deck"
    p.font.name = TITLE_FONT
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]

    sources = [
        "The Ultimate Guide to VMware Migration and Modernization on Azure, Microsoft eBook, 2025.",
        "Microsoft Learn: Start here to migrate from VMware to Azure.",
        "Microsoft Learn: Build a business case with Azure Migrate.",
        "Microsoft Learn: Plan and analyze VMware migrations using Azure Copilot migration agent.",
        "Microsoft Learn / Cloud Adoption Framework: Migrate workloads for Azure VMware Solution.",
        "Microsoft Learn / Cloud Adoption Framework: Azure landing zone review for Azure VMware Solution.",
        "Microsoft Learn / Well-Architected Framework: Integrate an Azure VMware Solution workload with Azure landing zones.",
        "Microsoft Learn / Well-Architected Framework: Azure VMware Solution workloads overview.",
    ]

    add_bullets(
        slide,
        Inches(0.92),
        Inches(1.8),
        Inches(11.2),
        Inches(4.8),
        sources,
        font_size=17,
        color=COLORS["white"],
    )
    add_footer(
        slide,
        "Presentation created in azure_arc_vsphere_jumpstart for reuse and further customization",
    )


def build_deck():
    add_title_slide()
    add_audience_map_slide()
    add_exec_summary_slide()
    add_section_divider(
        "Business leaders",
        "Business outcome framing, portfolio choices, and economics for VMware-to-Azure decisions.",
        COLORS["orange"],
    )
    add_business_case_slide()
    add_choice_matrix_slide()
    add_section_divider(
        "Technical decision makers",
        "Discovery, landing-zone preparation, migration waves, and hybrid operations.",
        COLORS["teal"],
    )
    add_discovery_slide()
    add_architecture_slide()
    add_landing_zone_slide()
    add_waves_slide()
    add_arc_slide()
    add_vdi_slide()
    add_cost_slide()
    add_risks_slide()
    add_action_plan_slide()
    add_sources_slide()
    prs.save(OUTPUT)


if __name__ == "__main__":
    build_deck()
    print(f"Created {OUTPUT}")
