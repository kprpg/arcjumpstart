from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE_DIR = Path(__file__).parent
OUTPUT = BASE_DIR / "azure_migrate_vmware_sample_walkthrough.pptx"

TITLE_COLOR = RGBColor(22, 49, 76)
ACCENT_COLOR = RGBColor(0, 120, 170)
TEXT_COLOR = RGBColor(40, 40, 40)
BG_COLOR = RGBColor(245, 248, 251)


def add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_title(slide, title, subtitle=""):
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.4), Inches(11.6), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.05), Inches(11.2), Inches(0.5)
        )
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(13)
        sp.font.color.rgb = ACCENT_COLOR


def add_bullets(slide, items, left=0.8, top=1.7, width=11.0, height=4.8):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(12)


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(11.0), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.RIGHT


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = [
        (
            "Azure Migrate Sample Package Walkthrough",
            "Enterprise, healthcare, workshop, and dependency-discovery demo assets",
            [
                "This deck explains how the RVTools import samples, enriched planning workbooks, checklist exports, and dependency-discovery samples fit together.",
                "Audience: migration architects, application owners, platform teams, and workshop facilitators.",
            ],
        ),
        (
            "Import-Safe Starting Point",
            "Use the RVTools-style workbook for Azure Migrate file import",
            [
                "Start with azure_migrate_vmware_rvtools_strict_sample.xlsx when you want a clean portal import test.",
                "Use azure_migrate_vmware_enterprise_import_sample.xlsx or healthcare_payor_provider_vmware_import_sample.xlsx for scaled scenarios.",
                "Only the RVTools tabs should be treated as import candidates.",
            ],
        ),
        (
            "Enriched Planning Copies",
            "Preserve RVTools tabs and add application context outside the import path",
            [
                "The enriched workbooks add ServerApplicationMap, ApplicationTopology, ClusterRelationships, ApplicationDependencies, WorkloadGroupsForAzureMigrate, and AzureMigrateGroupChecklist.",
                "These sheets are designed for planning, migration-wave review, and architecture workshops.",
                "They intentionally avoid changing the original 11 RVTools tabs.",
            ],
        ),
        (
            "Portal Grouping Aids",
            "Checklist sheet plus CSV exports for manual Azure Migrate group creation",
            [
                "Use AzureMigrateGroupChecklist inside the enriched workbooks to see recommended group names, member counts, tiers, waves, and portal path.",
                "Use the enterprise_azure_migrate_group_checklist.csv and enterprise_azure_migrate_group_members.csv style exports for copy/paste-heavy workshops.",
                "Equivalent healthcare and workshop CSVs are generated alongside the workbooks.",
            ],
        ),
        (
            "Workshop-Focused Workbook",
            "Dedicated lens for SQL, SAP, and Citrix migration discussions",
            [
                "azure_migrate_vmware_enterprise_sql_sap_citrix_workshop.xlsx narrows the conversation to the hardest clustering and dependency patterns.",
                "WorkshopFocusSummary, WorkshopFocusClusters, and WorkshopFocusDependencies speed up technical reviews.",
                "Use this workbook when the audience needs to reason about HA, failover, and control-plane preservation.",
            ],
        ),
        (
            "Dependency Discovery Sample",
            "Synthetic Azure Migrate dependency-analysis style output with variation",
            [
                "azure_migrate_dependency_discovery_variations_sample.xlsx includes DiscoverySummary, ObservedProcesses, ObservedConnections, and VariationGuide.",
                "The observed connections sample spans steady-state flows, nightly batch windows, bursty Citrix sessions, SQL AO heartbeats, healthcare DICOM traffic, and external service edges.",
                "Use the CSV sidecars when you want to filter, sort, or build dashboards around discovered dependencies.",
            ],
        ),
        (
            "Suggested Demo Flow",
            "How to present the sample set end-to-end",
            [
                "1. Import the RVTools workbook into Azure Migrate.",
                "2. Review the enriched copy to align applications, tiers, and migration waves.",
                "3. Use the checklist and group-member CSVs to create manual groups in the portal.",
                "4. Show the dependency-discovery sample to explain what native dependency analysis adds after discovery is enabled.",
            ],
        ),
    ]

    for index, (title, subtitle, bullets) in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide)
        add_title(slide, title, subtitle)
        add_bullets(slide, bullets)
        add_footer(slide, f"Azure Migrate sample walkthrough | Slide {index}")

    prs.save(OUTPUT)
    print(f"Created {OUTPUT}")


if __name__ == "__main__":
    build_deck()
