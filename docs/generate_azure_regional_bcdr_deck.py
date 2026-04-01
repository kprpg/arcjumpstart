from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

OUTPUT_PATH = r"c:\Users\gpillai\source\repos\arcjumpstart\docs\azure-regional-bcdr-slide-deck.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TITLE_COLOR = RGBColor(11, 44, 91)
ACCENT_COLOR = RGBColor(0, 120, 212)
TEXT_COLOR = RGBColor(33, 37, 41)
MUTED_COLOR = RGBColor(96, 103, 112)
BG_COLOR = RGBColor(248, 250, 252)

slides = [
    {
        "title": "Azure US East 2 Relocation and BCDR Strategy",
        "subtitle": "Planning for datacenter, regional, and wide-area East Coast disruption",
        "bullets": [],
    },
    {
        "title": "Why This Matters",
        "bullets": [
            "The scenario is not a routine outage or a standard DR exercise.",
            "The customer is planning for loss of US East 2 or a wider East Coast operating disruption.",
            "A US East 2-only resilience strategy is not sufficient for that threat model.",
        ],
    },
    {
        "title": "The Three Failure Scopes",
        "bullets": [
            "Single datacenter or availability-zone failure.",
            "Single-region failure in US East 2.",
            "Wider East Coast disruption affecting connectivity or the operating model.",
            "Each scenario requires a different control set.",
        ],
    },
    {
        "title": "Core Message",
        "bullets": [
            "Availability zones help with facility-level failures.",
            "Regional recovery patterns help with some region-level failures.",
            "Neither solves for loss of the broader East Coast operating geography.",
            "The appropriate hedge for that scenario is another geography.",
        ],
    },
    {
        "title": "Recommended Strategy",
        "bullets": [
            "Use multi-zone resilience inside the primary region.",
            "Use cross-region recovery for regional outages.",
            "Use cross-geography recovery outside the East Coast blast radius.",
        ],
    },
    {
        "title": "What Good Looks Like",
        "bullets": [
            "No single-instance production services.",
            "The disaster recovery environment exists before the crisis, not after.",
            "Data services replicate across regions.",
            "The application entry point supports controlled failover.",
            "Runbooks and recovery artifacts remain available outside the primary geography.",
        ],
    },
    {
        "title": "Region Options",
        "bullets": [
            "US East 2 plus Central US: default paired-region DR.",
            "US East 2 plus South Central US: lower-latency compromise.",
            "US East 2 plus West US 2 or West US 3: strongest geographic hedge.",
            "Central US primary plus US East 2 secondary: strong relocation option.",
        ],
    },
    {
        "title": "Recommendation by Scenario",
        "bullets": [
            "Use Central US as the default disaster recovery region for a balanced design.",
            "Use West US 2 or West US 3 when wider geographic separation is the priority.",
            "Evaluate Central US as a future primary if East Coast concentration needs to be reduced further.",
        ],
    },
    {
        "title": "Technical Pattern",
        "bullets": [
            "Azure Front Door or Traffic Manager for failover routing.",
            "Azure Site Recovery for VM-based workloads.",
            "Native geo-replication for PaaS data services.",
            "A secondary landing zone with networking, identity, logging, and secrets prebuilt.",
            "Immutable backups and tested restore procedures.",
        ],
    },
    {
        "title": "What Not to Assume",
        "bullets": [
            "Region pairing is not automatic DR.",
            "Platform-managed storage failover is not a full application recovery plan.",
            "Backups alone are not a business continuity strategy.",
            "A second east-side region does not fully hedge a wider East Coast disruption.",
        ],
    },
    {
        "title": "Decision Matrix",
        "table": {
            "columns": ["Option", "Position"],
            "rows": [
                ["US East 2 only", "Not sufficient"],
                ["US East 2 plus Central US", "Recommended default"],
                ["US East 2 plus South Central US", "Reasonable compromise"],
                [
                    "US East 2 plus West US 2 or West US 3",
                    "Recommended for stronger separation",
                ],
                [
                    "Central US primary plus US East 2 secondary",
                    "Strong relocation option",
                ],
            ],
        },
    },
    {
        "title": "Recommended Next Steps",
        "bullets": [
            "Confirm data residency constraints.",
            "Define recovery time and recovery point targets by workload tier.",
            "Select the target DR geography.",
            "Validate region service availability and quota.",
            "Build the secondary landing zone.",
            "Implement replication and failover.",
            "Run a real DR exercise.",
        ],
    },
    {
        "title": "Leadership Decision Required",
        "bullets": [
            "A balanced paired-region design centered on Central US.",
            "A stronger east-plus-west design for wider separation.",
            "A future primary-region shift if East Coast concentration needs to be reduced further.",
        ],
    },
    {
        "title": "Closing Statement",
        "bullets": [
            "If the organization is planning for the loss of US East 2 or a wider East Coast operating disruption, the right hedge is another geography, not just another zone in US East 2.",
        ],
    },
]


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_header_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.35)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()


def add_title(slide, title):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(12), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR


def add_footer(slide, index):
    box = slide.shapes.add_textbox(Inches(11.7), Inches(7.0), Inches(1.0), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(index)
    run.font.name = "Aptos"
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED_COLOR


def add_bullets(slide, bullets):
    box = slide.shapes.add_textbox(
        Inches(0.95), Inches(1.55), Inches(11.4), Inches(4.9)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for index, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(10)
        p.bullet = True


def add_subtitle(slide, subtitle):
    box = slide.shapes.add_textbox(Inches(0.95), Inches(1.8), Inches(11), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = subtitle
    run.font.name = "Aptos"
    run.font.size = Pt(24)
    run.font.color.rgb = MUTED_COLOR


def add_table(slide, columns, rows):
    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(columns), Inches(0.7), Inches(1.6), Inches(12.0), Inches(4.7)
    )
    table = table_shape.table
    for idx, name in enumerate(columns):
        cell = table.cell(0, idx)
        cell.text = name
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE_COLOR
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(16)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                RGBColor(255, 255, 255) if row_idx % 2 else RGBColor(240, 244, 248)
            )
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = TEXT_COLOR
    table.columns[0].width = Inches(5.2)
    table.columns[1].width = Inches(6.8)


for index, payload in enumerate(slides, start=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_bar(slide)
    add_title(slide, payload["title"])
    if payload.get("subtitle"):
        add_subtitle(slide, payload["subtitle"])
    if payload.get("bullets"):
        add_bullets(slide, payload["bullets"])
    if payload.get("table"):
        add_table(slide, payload["table"]["columns"], payload["table"]["rows"])
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.7),
        Inches(1.25),
        Inches(1.4),
        Inches(0.08),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_COLOR
    accent.line.fill.background()
    add_footer(slide, index)

prs.save(OUTPUT_PATH)
print(OUTPUT_PATH)
